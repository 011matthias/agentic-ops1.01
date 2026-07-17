"""Round-trip tests for tools/pptx_slide_ops.py (deck-foundation-v2 engine).

The fixture pptx is generated in-test with python-pptx (no binary in git).
Skips cleanly when python-pptx is unavailable (CI installs it via
`--with python-pptx` in ci.yml; a bare `uv run --with pytest pytest` skips).
"""

from __future__ import annotations

import importlib.util
import io
import sys
import zlib
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

REPO = Path(__file__).resolve().parents[2]
_spec = importlib.util.spec_from_file_location("pptx_slide_ops", REPO / "tools" / "pptx_slide_ops.py")
ops = importlib.util.module_from_spec(_spec)
sys.modules["pptx_slide_ops"] = ops
_spec.loader.exec_module(ops)


def _png_bytes(rgb: tuple[int, int, int]) -> bytes:
    """Minimal valid 1x1 PNG of the given color."""
    def chunk(tag: bytes, data: bytes) -> bytes:
        raw = tag + data
        return len(data).to_bytes(4, "big") + raw + zlib.crc32(raw).to_bytes(4, "big")

    ihdr = (1).to_bytes(4, "big") + (1).to_bytes(4, "big") + bytes([8, 2, 0, 0, 0])
    idat = zlib.compress(bytes([0, *rgb]))
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", idat)
        + chunk(b"IEND", b"")
    )


@pytest.fixture()
def deck(tmp_path):
    """3-slide fixture: s0 title, s1 picture + named textbox, s2 textbox
    with a literal 'N / M' footer run."""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    s0.shapes.title.text = "Alpha"

    s1 = prs.slides.add_slide(blank)
    png = tmp_path / "red.png"
    png.write_bytes(_png_bytes((255, 0, 0)))
    pic = s1.shapes.add_picture(str(png), Inches(1), Inches(1), Inches(2), Inches(2))
    pic.name = "hero-image"
    tb = s1.shapes.add_textbox(Inches(1), Inches(4), Inches(6), Inches(2))
    tb.name = "role-body"
    tb.text_frame.text = "First line"
    p2 = tb.text_frame.add_paragraph()
    p2.text = "Second line"
    p2.runs[0].font.size = Pt(10)
    p2.runs[0].font.bold = True

    s2 = prs.slides.add_slide(blank)
    ft = s2.shapes.add_textbox(Inches(11), Inches(7), Inches(1), Inches(0.4))
    ft.name = "footer-num"
    ft.text_frame.text = "9 / 9"

    path = tmp_path / "fixture.pptx"
    prs.save(str(path))
    return path


def _reload(path):
    return Presentation(str(path))


def test_duplicate_carries_picture_and_valid_rids(deck, tmp_path):
    prs = _reload(deck)
    new = ops.duplicate_slide(prs, 1)
    assert ops.find_shape(new, "hero-image") is not None
    assert ops.find_shape(new, "role-body").text_frame.text.startswith("First line")
    out = tmp_path / "dup.pptx"
    prs.save(str(out))
    prs2 = _reload(out)
    assert len(prs2.slides) == 4
    # the clone's picture must render from a resolvable image part
    clone = prs2.slides[3]
    pic = ops.find_shape(clone, "hero-image")
    assert pic.image.blob[:4] == b"\x89PNG"
    assert ops.validate_rids(out) == []


def test_duplicate_twice_unique_rids(deck, tmp_path):
    prs = _reload(deck)
    ops.duplicate_slide(prs, 1)
    ops.duplicate_slide(prs, 1)
    out = tmp_path / "dup2.pptx"
    prs.save(str(out))
    assert len(_reload(out).slides) == 5
    assert ops.validate_rids(out) == []


def test_delete_and_move(deck, tmp_path):
    prs = _reload(deck)
    ops.delete_slide(prs, 1)
    ops.move_slide(prs, 1, 0)  # footer slide to front
    out = tmp_path / "delmove.pptx"
    prs.save(str(out))
    prs2 = _reload(out)
    assert len(prs2.slides) == 2
    assert ops.find_shape(prs2.slides[0], "footer-num", required=False) is not None
    assert prs2.slides[1].shapes.title.text == "Alpha"
    assert ops.validate_rids(out) == []


def test_fill_paragraphs_grow_shrink_preserves_format(deck, tmp_path):
    prs = _reload(deck)
    shape = ops.find_shape(prs.slides[1], "role-body")
    ops.fill_paragraphs(shape, ["one", "two", "three", "four"])
    out = tmp_path / "fill.pptx"
    prs.save(str(out))
    shape2 = ops.find_shape(_reload(out).slides[1], "role-body")
    texts = [p.text for p in shape2.text_frame.paragraphs]
    assert texts == ["one", "two", "three", "four"]
    # grown paragraphs inherit the LAST donor paragraph's format (bold 10pt)
    assert shape2.text_frame.paragraphs[3].runs[0].font.bold is True

    ops.fill_paragraphs(shape2, ["only"])
    assert [p.text for p in shape2.text_frame.paragraphs] == ["only"]


def test_fill_outline_levels(deck, tmp_path):
    from pptx.oxml.ns import qn as _qn

    prs = _reload(deck)
    shape = ops.find_shape(prs.slides[1], "role-body")
    # give the fixture two levels: para 2 becomes lvl=1
    p2 = shape.text_frame.paragraphs[1]._p
    ppr = p2.get_or_add_pPr()
    ppr.set("lvl", "1")

    ops.fill_outline(shape, [(0, "Head A"), (1, "bullet 1"), (1, "bullet 2"), (0, ""), (0, "Head B"), (1, "bullet 3")])
    out = tmp_path / "outline.pptx"
    prs.save(str(out))
    shape2 = ops.find_shape(_reload(out).slides[1], "role-body")
    paras = shape2.text_frame.paragraphs
    assert [p.text for p in paras] == ["Head A", "bullet 1", "bullet 2", "", "Head B", "bullet 3"]
    lvls = [int(p._p.find(_qn("a:pPr")).get("lvl", "0")) if p._p.find(_qn("a:pPr")) is not None else 0 for p in paras]
    assert lvls == [0, 1, 1, 0, 0, 1]
    # bullets inherit the lvl-1 template's formatting (bold 10pt in fixture)
    assert paras[1].runs[0].font.bold is True


def test_patch_runs_exact_and_substring(deck):
    prs = _reload(deck)
    shape = ops.find_shape(prs.slides[1], "role-body")
    assert ops.patch_runs(shape, {"First line": "Patched"}) == 1
    assert shape.text_frame.paragraphs[0].runs[0].text == "Patched"
    assert ops.patch_runs(shape, {"Patch": "Match"}, substring=True) == 1
    assert shape.text_frame.paragraphs[0].runs[0].text == "Matched"


def test_swap_picture_changes_blob(deck, tmp_path):
    prs = _reload(deck)
    blue = tmp_path / "blue.png"
    blue.write_bytes(_png_bytes((0, 0, 255)))
    ops.swap_picture(prs.slides[1], "hero-image", blue)
    out = tmp_path / "swap.pptx"
    prs.save(str(out))
    pic = ops.find_shape(_reload(out).slides[1], "hero-image")
    assert pic.image.blob == _png_bytes((0, 0, 255))
    assert ops.validate_rids(out) == []


def test_repoint_picture_shares_part(deck, tmp_path):
    prs = _reload(deck)
    clone = ops.duplicate_slide(prs, 1)
    blue = tmp_path / "blue2.png"
    blue.write_bytes(_png_bytes((0, 0, 255)))
    ops.swap_picture(prs.slides[1], "hero-image", blue)
    # clone still red; repoint it at the (now blue) original
    ops.repoint_picture(clone, "hero-image", prs.slides[1], "hero-image")
    out = tmp_path / "repoint.pptx"
    prs.save(str(out))
    prs2 = _reload(out)
    assert ops.find_shape(prs2.slides[3], "hero-image").image.blob == _png_bytes((0, 0, 255))
    assert ops.validate_rids(out) == []


def test_hidden_flag_and_footer_renumber(deck, tmp_path):
    prs = _reload(deck)
    ops.set_hidden(prs.slides[0], True)
    assert ops.is_hidden(prs.slides[0])
    # footer counts only VISIBLE slides: fixture has 3, one hidden -> 2,
    # and the footer slide is visible position 2
    assert ops.renumber_footers(prs) == 1
    ft = ops.find_shape(prs.slides[2], "footer-num")
    assert ft.text_frame.text == "2 / 2"
    ops.set_hidden(prs.slides[0], False)
    assert not ops.is_hidden(prs.slides[0])


def test_roster_reports_layout_and_hidden(deck):
    prs = _reload(deck)
    ops.set_hidden(prs.slides[2], True)
    tmp = deck.parent / "roster.pptx"
    prs.save(str(tmp))
    rows = ops.roster(tmp)
    assert len(rows) == 3
    assert rows[0]["head"] == "Alpha"
    assert rows[2]["hidden"] is True
