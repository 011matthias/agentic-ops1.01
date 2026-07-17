# /// script
# dependencies = ["python-pptx>=1.0.2"]
# ///
"""pptx surgery primitives: duplicate/delete/reorder slides, patch text and
pictures, SmartArt text, footer renumbering, embedded-font + rId integrity
checks.

Built for the Brisken deck-foundation-v2 system (clone-and-patch on a
library pptx derived from Dirk's reference deck), but deck-agnostic: every
function takes a Presentation or a path. Importable as a module
(`import pptx_slide_ops`) and runnable as a CLI:

    uv run tools/pptx_slide_ops.py roster <file.pptx>
    uv run tools/pptx_slide_ops.py validate <file.pptx>

Design constraints honored here:
- Slide duplication is IN-PACKAGE only (same .pptx). Cross-package cloning
  is deliberately unsupported; composing decks = copy the whole library
  file, then duplicate/delete/reorder inside it. This keeps every layout /
  media / font relationship valid with zero remapping beyond the slide's
  own rels.
- Text patching is run-preserving: formatting always comes from an existing
  run/paragraph of the donor slide, never re-specified in code.
- SmartArt is patched in place, never duplicated (shared diagram parts).
"""

from __future__ import annotations

import argparse
import copy
import re
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

# Slide-part rels that must NOT be copied onto a duplicate: the layout rel
# already exists on the fresh slide; notes are dropped by design.
_DUP_SKIP_RELTYPES = {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE}

DIAGRAM_DATA_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"
)


# ---------------------------------------------------------------- shapes

def iter_shapes(container, recursive: bool = True):
    """Yield shapes; descend into groups when recursive."""
    for shp in container.shapes:
        yield shp
        if recursive and shp.shape_type is not None and shp.shape_type == 6:  # GROUP
            yield from _iter_group(shp)


def _iter_group(group):
    for shp in group.shapes:
        yield shp
        if shp.shape_type == 6:
            yield from _iter_group(shp)


def find_shape(slide, name: str, required: bool = True):
    """First shape (recursing groups) whose .name equals `name`."""
    for shp in iter_shapes(slide):
        if shp.name == name:
            return shp
    if required:
        raise KeyError(f"shape {name!r} not found on slide")
    return None


def rename_shape(slide, old: str, new: str) -> None:
    find_shape(slide, old).name = new


# ---------------------------------------------------------------- slides

def slide_index(prs, slide) -> int:
    for i, s in enumerate(prs.slides):
        if s.slide_id == slide.slide_id:
            return i
    raise ValueError("slide not in presentation")


def has_smartart(slide) -> bool:
    for shp in slide.shapes:
        el = shp._element
        for gd in el.iter(qn("a:graphicData")):
            if "diagram" in (gd.get("uri") or ""):
                return True
    return False


def duplicate_slide(prs, index: int, allow_smartart: bool = False):
    """Clone slide `index` inside the same package; returns the new slide
    (appended last — use move_slide to position it).

    The clone gets its own slide part; every relationship of the source
    slide part (images, charts, hyperlinks) is re-created on the clone and
    every r:-namespace attribute in the copied shape tree is remapped to
    the new rIds, so duplicates never share or collide rIds.
    """
    src = prs.slides[index]
    if has_smartart(src) and not allow_smartart:
        raise ValueError(
            "refusing to duplicate a SmartArt slide (diagram parts would be "
            "shared); patch it in place instead"
        )
    new = prs.slides.add_slide(src.slide_layout)

    # strip the fresh slide's layout-inherited placeholders
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)

    # background
    src_bg = src._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if src_bg is not None:
        new._element.find(qn("p:cSld")).insert(0, copy.deepcopy(src_bg))

    # hidden flag
    show = src._element.get("show")
    if show is not None:
        new._element.set("show", show)

    # re-create rels; build old-rId -> new-rId map
    rid_map: dict[str, str] = {}
    for rid, rel in src.part.rels.items():
        if rel.reltype in _DUP_SKIP_RELTYPES:
            continue
        if rel.is_external:
            new_rid = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rid] = new_rid

    # copy shapes, remapping every r:* attribute through rid_map
    src_sptree = src._element.find(qn("p:cSld") + "/" + qn("p:spTree"))
    new_sptree = new._element.find(qn("p:cSld") + "/" + qn("p:spTree"))
    skip = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in src_sptree:
        if child.tag in skip:
            continue
        clone = copy.deepcopy(child)
        _remap_r_attrs(clone, rid_map)
        new_sptree.append(clone)
    return new


def _remap_r_attrs(el, rid_map: dict[str, str]) -> None:
    for node in el.iter():
        for attr, val in list(node.attrib.items()):
            if attr.startswith("{" + R_NS + "}") and val in rid_map:
                node.set(attr, rid_map[val])


def delete_slide(prs, index: int) -> None:
    """Drop slide from sldIdLst + drop its rel (part becomes unreachable
    and is excluded at save — the _rebuild_dcw_xml pattern)."""
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    el = sld_ids[index]
    prs.part.drop_rel(el.get(qn("r:id")))
    sld_id_lst.remove(el)


def move_slide(prs, from_index: int, to_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    el = sld_ids[from_index]
    sld_id_lst.remove(el)
    sld_id_lst.insert(to_index, el)


def set_hidden(slide, hidden: bool) -> None:
    if hidden:
        slide._element.set("show", "0")
    elif slide._element.get("show") is not None:
        del slide._element.attrib["show"]


def is_hidden(slide) -> bool:
    return slide._element.get("show") == "0"


# ---------------------------------------------------------------- text

def patch_runs(shape, mapping: dict[str, str], substring: bool = False) -> int:
    """Replace run texts. Exact (stripped) match by default; substring mode
    replaces within runs. Formatting untouched. Returns replacement count."""
    n = 0
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if substring:
                new = run.text
                for old, repl in mapping.items():
                    if old in new:
                        new = new.replace(old, repl)
                if new != run.text:
                    run.text = new
                    n += 1
            else:
                key = run.text.strip()
                if key in mapping:
                    run.text = mapping[key]
                    n += 1
    return n


def fill_paragraphs(shape, texts: list[str]) -> None:
    """Paragraph-level replacement preserving formatting: paragraph i keeps
    its own format, its first run gets texts[i], surplus runs are removed.
    Extra texts clone the LAST donor paragraph; surplus paragraphs are
    deleted. Donor must have at least one paragraph with one run."""
    tf = shape.text_frame
    paras = tf.paragraphs
    if not paras or not any(p.runs for p in paras):
        raise ValueError(f"shape {shape.name!r} has no run to inherit format from")

    txbody = tf._txBody
    p_els = txbody.findall(qn("a:p"))

    # grow: clone last paragraph element
    while len(p_els) < len(texts):
        clone = copy.deepcopy(p_els[-1])
        p_els[-1].addnext(clone)
        p_els = txbody.findall(qn("a:p"))
    # shrink
    while len(p_els) > len(texts):
        txbody.remove(p_els[-1])
        p_els = txbody.findall(qn("a:p"))

    for p_el, text in zip(p_els, texts):
        runs = p_el.findall(qn("a:r"))
        if not runs:
            # paragraph with no run (blank spacer): borrow a run from a donor
            donor = next(p for p in txbody.findall(qn("a:p")) if p.findall(qn("a:r")))
            runs = [copy.deepcopy(donor.findall(qn("a:r"))[0])]
            p_el.append(runs[0])
        for surplus in runs[1:]:
            p_el.remove(surplus)
        t = runs[0].find(qn("a:t"))
        if t is None:
            t = etree.SubElement(runs[0], qn("a:t"))
        t.text = text


def fill_outline(shape, items: list[tuple[int, str]]) -> None:
    """Rebuild a leveled text body (head/bullet outlines) from donor styles.

    items = [(level, text), ...]. For each level present in the donor, the
    FIRST donor paragraph at that level is the style template; every output
    paragraph is a clone of its level's template with the text swapped in.
    Empty text keeps the paragraph as a spacer (runs removed)."""
    tf = shape.text_frame
    txbody = tf._txBody
    p_els = txbody.findall(qn("a:p"))
    if not p_els:
        raise ValueError(f"shape {shape.name!r} has no paragraphs")

    templates: dict[int, etree._Element] = {}
    for p_el in p_els:
        ppr = p_el.find(qn("a:pPr"))
        lvl = int(ppr.get("lvl", "0")) if ppr is not None else 0
        if lvl not in templates and p_el.findall(qn("a:r")):
            templates[lvl] = copy.deepcopy(p_el)

    missing = {lvl for lvl, _ in items} - set(templates)
    if missing:
        raise ValueError(f"donor {shape.name!r} has no template for level(s) {sorted(missing)}")

    for p_el in p_els:
        txbody.remove(p_el)

    for lvl, text in items:
        p_el = copy.deepcopy(templates[lvl])
        runs = p_el.findall(qn("a:r"))
        for surplus in runs[1:]:
            p_el.remove(surplus)
        if text:
            t = runs[0].find(qn("a:t"))
            if t is None:
                t = etree.SubElement(runs[0], qn("a:t"))
            t.text = text
        else:
            p_el.remove(runs[0])
        txbody.append(p_el)


def swap_picture(slide, shape_name: str, image_path: str | Path):
    """Point the named Picture at a new image part (added to the package);
    other slides sharing the old image part are unaffected."""
    pic = find_shape(slide, shape_name)
    image_part, new_rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic._element.find(".//" + qn("a:blip"))
    if blip is None:
        raise ValueError(f"shape {shape_name!r} has no a:blip (not a picture?)")
    blip.set(qn("r:embed"), new_rid)
    return image_part


def repoint_picture(slide, shape_name: str, src_slide, src_shape_name: str) -> None:
    """Point the named Picture at the image part another slide's Picture
    already uses (same package). No new media part is created."""
    src_pic = find_shape(src_slide, src_shape_name)
    src_blip = src_pic._element.find(".//" + qn("a:blip"))
    if src_blip is None:
        raise ValueError(f"source shape {src_shape_name!r} has no a:blip")
    image_part = src_slide.part.related_part(src_blip.get(qn("r:embed")))
    rid = slide.part.relate_to(image_part, RT.IMAGE)
    blip = find_shape(slide, shape_name)._element.find(".//" + qn("a:blip"))
    if blip is None:
        raise ValueError(f"target shape {shape_name!r} has no a:blip")
    blip.set(qn("r:embed"), rid)


def smartart_set_text(slide, mapping: dict[str, str]) -> int:
    """Exact-match replace of <a:t> texts inside the slide's SmartArt
    diagram data part(s). In-place; returns replacement count."""
    n = 0
    for rid, rel in slide.part.rels.items():
        if rel.reltype != DIAGRAM_DATA_RELTYPE:
            continue
        root = rel.target_part._element
        for t in root.iter(qn("a:t")):
            if t.text and t.text.strip() in mapping:
                t.text = mapping[t.text.strip()]
                n += 1
    return n


_FOOTER_RE = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")


def renumber_footers(prs) -> int:
    """Rewrite literal 'N / M' footer runs to visible-position / total.
    Slide-number placeholder fields renumber themselves; this only touches
    literal text runs. Returns rewrite count."""
    visible = [s for s in prs.slides if not is_hidden(s)]
    total = len(visible)
    n = 0
    for pos, slide in enumerate(visible, 1):
        for shp in iter_shapes(slide):
            if not shp.has_text_frame:
                continue
            for para in shp.text_frame.paragraphs:
                for run in para.runs:
                    if _FOOTER_RE.match(run.text or ""):
                        run.text = f"{pos} / {total}"
                        n += 1
    return n


# ---------------------------------------------------------------- checks

def check_embedded_fonts(path: str | Path) -> dict:
    """The reference embeds Poppins/Open Sans; a valid build must keep the
    font parts AND the embeddedFontLst declaration."""
    with zipfile.ZipFile(path) as z:
        fntdata = [n for n in z.namelist() if n.startswith("ppt/fonts/")]
        pres = z.read("ppt/presentation.xml").decode("utf-8", "ignore")
    return {
        "font_parts": len(fntdata),
        "declared": "embeddedFontLst" in pres,
        "ok": bool(fntdata) and "embeddedFontLst" in pres,
    }


def validate_rids(path: str | Path) -> list[str]:
    """Every r:* attribute referenced in a slide's XML must resolve to a
    relationship in that slide's .rels. Returns violation strings."""
    problems = []
    with zipfile.ZipFile(path) as z:
        slides = sorted(
            n for n in z.namelist()
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", n)
        )
        for name in slides:
            rels_name = f"ppt/slides/_rels/{name.rsplit('/', 1)[1]}.rels"
            defined = set()
            if rels_name in z.namelist():
                rroot = etree.fromstring(z.read(rels_name))
                defined = {r.get("Id") for r in rroot}
            root = etree.fromstring(z.read(name))
            used = {
                val
                for node in root.iter()
                for attr, val in node.attrib.items()
                if attr.startswith("{" + R_NS + "}")
            }
            for rid in sorted(used - defined):
                problems.append(f"{name}: {rid} referenced but not defined")
    return problems


def roster(path: str | Path) -> list[dict]:
    prs = Presentation(str(path))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        head = ""
        for shp in iter_shapes(slide):
            if shp.has_text_frame and shp.text_frame.text.strip():
                head = shp.text_frame.text.strip().splitlines()[0][:60]
                break
        out.append({
            "n": i,
            "layout": slide.slide_layout.name,
            "hidden": is_hidden(slide),
            "shapes": len(slide.shapes),
            "head": head,
        })
    return out


# ---------------------------------------------------------------- CLI

def _main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    sub = ap.add_subparsers(dest="cmd", required=True)
    for cmd in ("roster", "validate", "fonts"):
        p = sub.add_parser(cmd)
        p.add_argument("pptx")
    args = ap.parse_args(argv)

    if args.cmd == "roster":
        for row in roster(args.pptx):
            flag = " HIDDEN" if row["hidden"] else ""
            print(f"{row['n']:>3}  [{row['layout']}] shapes={row['shapes']}{flag}  {row['head']}")
        return 0
    if args.cmd == "validate":
        problems = validate_rids(args.pptx)
        for p in problems:
            print("RID-FAIL", p)
        print("rIds:", "OK" if not problems else f"{len(problems)} problem(s)")
        return 1 if problems else 0
    if args.cmd == "fonts":
        res = check_embedded_fonts(args.pptx)
        print(f"font parts: {res['font_parts']}  declared: {res['declared']}  -> {'OK' if res['ok'] else 'MISSING'}")
        return 0 if res["ok"] else 1
    return 2


if __name__ == "__main__":
    raise SystemExit(_main(sys.argv[1:]))
