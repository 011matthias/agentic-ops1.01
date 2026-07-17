# /// script
# dependencies = ["python-pptx>=1.0.2", "pyyaml"]
# ///
"""Compose a deck from a spec: copy the library, duplicate/patch/reorder/
delete donor slides, verify, and run the banned-terms gate.

    uv run compose.py mdh-commodities            # specs/mdh-commodities.yaml

Spec format (per slide: a pattern from library-manifest.json + patch
directives keyed by semantic role name):

    deck: mdh-commodities
    output: "Brisken - Market Data Hub Commodities 2026-08 PROPOSAL"
    slides:
      - pattern: cover
        patch:
          cover-title: {runs: {"Solutions Overview": "Market Data Hub for Commodities"}}
      - pattern: usecase
        patch:
          uc-title: {set: "Commodity Price Curves"}
          uc-col1-body: {outline: [[0, "Description"], [1, "..."], ...]}
      - pattern: platform-context
        patch:
          pc-app-mdh: {highlight: true}

Directives: set (single para), paras (list), outline ([[lvl, text], ...]),
runs (exact run-text map — 0 hits = hard fail, the drift tripwire),
sub (substring map), image (swap picture blob), highlight (cyan fill +
white bold text), smartart (exact a:t map on the slide's diagram parts).

Output: {main}/.scratch/deckgen-v2/dist/{deck}/{output}.pptx, then
tools/validate-demo-material.py --client brisken over that folder.
"""

from __future__ import annotations

import importlib.util
import json
import shutil
import subprocess
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor

_here = Path(__file__).parent


def _load(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


common = _load("common", _here / "common.py")
ops = _load("pptx_slide_ops", Path(__file__).resolve().parents[6] / "tools" / "pptx_slide_ops.py")

CYAN = RGBColor(0x00, 0xB5, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MANIFEST = json.loads((_here / "library-manifest.json").read_text(encoding="utf-8"))["patterns"]


def apply_directive(prs, slide, role: str, directive: dict, spec_dir: Path) -> None:
    shape = None if role == "@smartart" else ops.find_shape(slide, role)
    for key, val in directive.items():
        if key == "set":
            ops.fill_paragraphs(shape, [val])
        elif key == "paras":
            ops.fill_paragraphs(shape, list(val))
        elif key == "outline":
            ops.fill_outline(shape, [(int(l), t) for l, t in val])
        elif key == "runs":
            hits = ops.patch_runs(shape, dict(val))
            if hits < len(val):
                raise SystemExit(
                    f"runs patch on {role!r}: only {hits}/{len(val)} donor texts "
                    f"matched — donor content drifted, update the spec"
                )
        elif key == "sub":
            ops.patch_runs(shape, dict(val), substring=True)
        elif key == "image":
            p = Path(val)
            ops.swap_picture(slide, role, p if p.is_absolute() else spec_dir / p)
        elif key == "from_shape":
            src_pattern, src_role = val
            src_idx = MANIFEST[src_pattern]["donor"] - 1
            ops.repoint_picture(slide, role, prs.slides[src_idx], src_role)
        elif key == "highlight":
            if val:
                shape.fill.solid()
                shape.fill.fore_color.rgb = CYAN
                for para in shape.text_frame.paragraphs:
                    for r in para.runs:
                        r.font.color.rgb = WHITE
                        r.font.bold = True
        elif key == "smartart":
            n = ops.smartart_set_text(slide, dict(val))
            if n < len(val):
                raise SystemExit(f"smartart patch: only {n}/{len(val)} texts matched")
        else:
            raise SystemExit(f"unknown directive {key!r} on role {role!r}")


def main() -> int:
    if len(sys.argv) != 2:
        print(__doc__)
        return 2
    spec_path = _here / "specs" / f"{sys.argv[1]}.yaml"
    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    manifest = MANIFEST

    # validate spec against manifest before touching anything
    for i, sl in enumerate(spec["slides"], 1):
        pat = sl.get("pattern")
        if pat not in manifest:
            raise SystemExit(f"slide {i}: unknown pattern {pat!r}")
        for role in sl.get("patch", {}):
            if role != "@smartart" and role not in manifest[pat]["roles"]:
                raise SystemExit(
                    f"slide {i} ({pat}): role {role!r} not in manifest roles "
                    f"{manifest[pat]['roles']}"
                )

    out_dir = common.dist_dir() / spec["deck"]
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / (spec["output"] + ".pptx")
    shutil.copyfile(common.library_path(), out_path)
    prs = Presentation(str(out_path))

    # phase 1 — resolve/duplicate donors (append-only; donor indices stable)
    used: dict[int, int] = {}
    slide_objs = []
    for sl in spec["slides"]:
        donor_idx = manifest[sl["pattern"]]["donor"] - 1
        used[donor_idx] = used.get(donor_idx, 0) + 1
        if used[donor_idx] == 1:
            slide_objs.append(prs.slides[donor_idx])
        else:
            slide_objs.append(ops.duplicate_slide(prs, donor_idx))

    # phase 2 — patch
    for sl, slide in zip(spec["slides"], slide_objs):
        for role, directive in sl.get("patch", {}).items():
            apply_directive(prs, slide, role, directive, spec_path.parent)

    # phase 3 — order to spec
    for target, slide in enumerate(slide_objs):
        ops.move_slide(prs, ops.slide_index(prs, slide), target)

    # phase 4 — drop everything else
    while len(prs.slides) > len(slide_objs):
        ops.delete_slide(prs, len(slide_objs))

    # phase 5 — housekeeping + save
    for slide in prs.slides:
        ops.set_hidden(slide, False)
    ops.renumber_footers(prs)
    prs.save(str(out_path))

    # verify structure
    prs2 = Presentation(str(out_path))
    assert len(prs2.slides) == len(spec["slides"]), "slide count mismatch"
    assert not any(ops.is_hidden(s) for s in prs2.slides), "hidden slide leaked"
    fonts = ops.check_embedded_fonts(out_path)
    rids = ops.validate_rids(out_path)
    if not fonts["ok"] or rids:
        print("STRUCT FAIL fonts:", fonts, "rids:", rids)
        return 1
    print(f"composed: {out_path}  ({len(prs2.slides)} slides; fonts {fonts['font_parts']} parts; rIds OK)")

    # banned-terms gate (B2): validate the deck folder
    validator = common.main_clone_root() / "tools" / "validate-demo-material.py"
    res = subprocess.run(
        ["uv", "run", str(validator), "--client", "brisken", "--dir", str(out_dir)],
        capture_output=True, text=True,
    )
    print(res.stdout.strip())
    if res.returncode != 0:
        print(res.stderr.strip())
        print("BANNED-TERMS GATE FAILED")
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
