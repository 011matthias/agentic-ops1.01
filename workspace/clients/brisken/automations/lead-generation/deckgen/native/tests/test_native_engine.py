"""Engine tests for the native deck family. Run without the gitignored
reference binary: placeholder assets are generated per test via Pillow.

    uv run --no-project --with pytest --with python-pptx --with pillow \
        --with pyyaml pytest -q native/tests
"""
import importlib.util
import re
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest
import yaml

NATIVE = Path(__file__).resolve().parents[1]
REPO = NATIVE.parents[6]


def _load(name, path):
    spec = importlib.util.spec_from_file_location(name, path)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


T = _load("tokens", NATIVE / "tokens.py")
logosets = _load("logosets", NATIVE / "logosets.py")
grammar = _load("grammar", NATIVE / "grammar.py")
assets_mod = _load("assets", NATIVE / "assets.py")
compose = _load("compose", NATIVE / "compose.py")

PRODUCT_DECKS = ["market-data-hub", "mdh-commodities", "smart-trading", "digital-co-worker"]


@pytest.fixture()
def dummy_assets(tmp_path):
    from PIL import Image
    a = tmp_path / "assets"
    (a / "logos").mkdir(parents=True)
    for name in ("brisken_dark.png", "brisken_reversed.png", "sap_certified.png"):
        Image.new("RGBA", (120, 30), (10, 10, 10, 255)).save(a / name)
    for name in logosets.ALL_WALL_LOGOS:
        Image.new("RGBA", (80, 40), (40, 40, 40, 255)).save(a / "logos" / f"{name}.png")
    return a


def _mini_spec(slides, palette="overview", output="test-deck"):
    return {"deck": "test", "output": output, "palette": palette, "slides": slides}


def _slide_texts(pptx: Path):
    out = {}
    with zipfile.ZipFile(pptx) as z:
        for n in z.namelist():
            m = re.match(r"ppt/slides/slide(\d+)\.xml$", n)
            if m:
                xml = z.read(n).decode("utf8", "ignore")
                out[int(m.group(1))] = " ".join(re.findall(r"<a:t>(.*?)</a:t>", xml, re.S))
    return out


ALL_TYPES_SPEC = [
    {"type": "cover", "eyebrow": "Powered by OnePilot", "title": "Test", "sub": "Solutions",
     "tagline": "A tagline."},
    {"type": "divider", "num": "01", "title": "Section", "sub": "Sub."},
    {"type": "short_version", "head": "Head.", "cards": [
        {"label": "A", "body": "a"}, {"label": "B", "body": "b"},
        {"label": "C", "body": "c"}, {"label": "D", "body": "d"}], "note": "Note."},
    {"type": "about", "kick": "About", "head": "Head.", "intro": "Intro.", "cards": [
        {"label": "A", "body": "a"}, {"label": "B", "body": "b"}, {"label": "C", "body": "c"}],
     "footnote": "Foot."},
    {"type": "customers", "head": "Head.", "caption": "Caption."},
    {"type": "hub", "kick": "Hub", "head": "Head.", "intro": "Intro.", "cards": [
        {"label": "A", "body": "a"}, {"label": "B", "body": "b"}, {"label": "C", "body": "c"}]},
    {"type": "hierarchy", "focal_app": "mdh"},
    {"type": "functional", "kick": "Func", "head": "Head.", "product": "Prod",
     "sources": ["s1", "s2"], "dests": ["d1", "d2"], "note": "Note.", "badge": True},
    {"type": "governance", "kick": "Gov", "head": "Head.", "intro": "Intro.", "cards": [
        {"label": "A", "body": "a"}, {"label": "B", "body": "b"},
        {"label": "C", "body": "c"}, {"label": "D", "body": "d"}], "footline": "Foot."},
    {"type": "app", "kick": "App", "head": "Head.", "gloss": "Gloss.",
     "problem": ["p1", "p2"], "steps": ["s1", "s2", "s3"], "freed": ["f1", "f2"],
     "connects": "CONNECTS TO x"},
    {"type": "usecase", "kick": "UC", "title": "Title", "gloss": "Gloss.", "before": "Before.",
     "steps": ["s1", "s2", "s3"], "checkpoint": "A person checks."},
    {"type": "problem", "kick": "Problem", "head": "Head.", "bands": [
        {"label": "Your sources", "body": "b1"},
        {"label": "The manual middle", "body": "b2"},
        {"label": "Your systems", "body": "b3"}], "stat": "Stat (Named Source)."},
    {"type": "success", "head": "Head.", "intro": "Intro.", "stories": [
        {"industry": "X", "deployment": "SAP", "what": "w", "highlight": "h",
         "replaced": "Replaced a manual solution."},
        {"industry": "Y", "deployment": "SAP", "what": "w", "highlight": "h",
         "replaced": "Replaced a manual solution."}]},
    {"type": "compare", "kick": "Cmp", "head": "Head.", "left_label": "L",
     "left_items": ["l1", "l2"], "right_label": "R", "right_items": ["r1", "r2"]},
    {"type": "contact"},
]


def test_palettes_registry():
    assert set(T.PALETTES) == {"overview"} | set(PRODUCT_DECKS)
    accents = [str(p.accent) for p in T.PALETTES.values()]
    assert len(accents) == len(set(accents)), "accents must be pairwise distinct"
    brights = [str(p.bright) for p in T.PALETTES.values()]
    assert len(brights) == len(set(brights))
    base = {str(c) for c in (T.INK, T.PAPER, T.NEUTRAL, T.MUTED, T.FAINT, T.LINE)}
    assert not (set(accents) & base)
    assert T.PALETTES["overview"].signature == "none"
    valid_sigs = {"none", "rail-left", "bar-right", "baseline", "corner-dots"}
    assert all(p.signature in valid_sigs for p in T.PALETTES.values())
    assert T.ALLOWED_FONTS == {"Century Gothic", "Segoe UI", "Segoe UI Semibold"}


def test_grammar_smoke_all_types(tmp_path, dummy_assets):
    spec = _mini_spec(ALL_TYPES_SPEC)
    assert compose.validate_spec(spec) == []
    dest = compose.build(spec, tmp_path, dummy_assets)
    from pptx import Presentation
    prs = Presentation(str(dest))
    assert len(prs.slides._sldIdLst) == len(ALL_TYPES_SPEC)
    assert prs.slide_width == int(T.SW * T.EMU)
    assert prs.slide_height == int(T.SH * T.EMU)


def test_footer_pagination(tmp_path, dummy_assets):
    spec = _mini_spec([
        ALL_TYPES_SPEC[0],           # cover (unnumbered)
        ALL_TYPES_SPEC[2],           # short_version -> 01
        ALL_TYPES_SPEC[1],           # divider (unnumbered)
        ALL_TYPES_SPEC[9],           # app -> 02
        ALL_TYPES_SPEC[10],          # usecase -> 03
        ALL_TYPES_SPEC[14],          # contact (unnumbered)
    ])
    dest = compose.build(spec, tmp_path, dummy_assets)
    texts = _slide_texts(dest)
    SIG = "TreasuryCentral, powered by OnePilot"
    assert re.search(r"\b01\b", texts[2]) and SIG in texts[2]
    assert re.search(r"\b02\b", texts[4]) and SIG in texts[4]
    assert re.search(r"\b03\b", texts[5]) and SIG in texts[5]
    for unnumbered in (1, 3, 6):
        assert SIG not in texts[unnumbered], \
            f"slide {unnumbered} must carry no footer (cover/divider/contact)"


def test_usecase_overflow_guard():
    sl = dict(ALL_TYPES_SPEC[10])
    sl["steps"] = [f"s{i}" for i in range(7)]
    errs = compose.validate_spec(_mini_spec([sl]))
    assert any("max 6" in e for e in errs)


def test_usecase_step_length_guard():
    sl = dict(ALL_TYPES_SPEC[10])
    sl["steps"] = ["x" * 130]
    errs = compose.validate_spec(_mini_spec([sl]))
    assert any("124 chars" in e for e in errs)


def test_problem_needs_three_bands():
    sl = dict(ALL_TYPES_SPEC[11])
    sl["bands"] = sl["bands"][:2]
    errs = compose.validate_spec(_mini_spec([sl]))
    assert any("3 bands" in e for e in errs)


def test_spec_validation_committed_specs():
    specs = sorted((NATIVE / "specs").glob("*.yaml"))
    assert specs, "no committed specs found"
    for p in specs:
        spec = yaml.safe_load(p.read_text(encoding="utf-8"))
        errs = compose.validate_spec(spec)
        assert errs == [], f"{p.name}: {errs}"


def test_em_dash_rejected_in_spec():
    sl = dict(ALL_TYPES_SPEC[2])
    sl = {**sl, "head": "A head — with an em-dash."}
    errs = compose.validate_spec(_mini_spec([sl]))
    assert any("em-dash" in e for e in errs)


def test_fonts_native_and_em_dash_zero(tmp_path, dummy_assets):
    dest = compose.build(_mini_spec(ALL_TYPES_SPEC), tmp_path, dummy_assets)
    with zipfile.ZipFile(dest) as z:
        names = z.namelist()
        assert not [n for n in names if n.startswith("ppt/fonts/")]
        faces = set()
        for n in names:
            if n.startswith("ppt/slides/") and n.endswith(".xml"):
                data = z.read(n)
                assert b"\xe2\x80\x94" not in data, f"em-dash in {n}"
                faces |= {f.decode() for f in re.findall(rb'typeface="([^"]+)"', data)}
    assert faces <= T.ALLOWED_FONTS, f"non-family fonts: {faces - T.ALLOWED_FONTS}"


def test_banned_terms_canary(tmp_path, dummy_assets):
    sl = dict(ALL_TYPES_SPEC[9])
    sl = {**sl, "gloss": "Runs on SAP BTP for speed."}
    dest = compose.build(_mini_spec([sl], output="canary"), tmp_path, dummy_assets)
    validator = REPO / "tools" / "validate-demo-material.py"
    res = subprocess.run(["uv", "run", str(validator), "--client", "brisken", str(dest)],
                         capture_output=True, text=True)
    assert res.returncode != 0, "validator must fail on BTP in a composed deck"
    assert "BTP" in res.stdout


def test_badge_identity_pins():
    assert assets_mod.BRAND_ASSETS["image12.png"] == (
        "sap_certified.png", "0a4e1f2466238d0723f83a4edb2d9475")
    assert "image34.png" not in assets_mod.BRAND_ASSETS, \
        "image34 is the Fortitude Re art, never the SAP badge"
    # The reference extraction still pins 20 logos (provenance/fallback). The
    # wall itself now renders the curated transparent library (logosets),
    # decoupled from the reference names since 2026-07-27 (Dirk's transparent
    # set: Ford/Siemens added, Angus transparent, Beautycounter/Global Brands
    # retired).
    assert len(assets_mod.CLIENT_LOGOS) == 20


def test_logo_sets_wellformed():
    assert logosets.DEFAULT_SET in logosets.LOGO_SETS
    assert len(logosets.LOGO_SETS["master"]) == 20
    # every set entry resolves to a curated-library name
    for name, members in logosets.LOGO_SETS.items():
        assert members, f"logo set {name} is empty"
        assert set(members) <= set(logosets.ALL_WALL_LOGOS)
    # ALL_WALL_LOGOS is exactly the union, sorted + deduped
    union = sorted({m for ms in logosets.LOGO_SETS.values() for m in ms})
    assert logosets.ALL_WALL_LOGOS == union


def test_customers_logo_set_validation(tmp_path, dummy_assets):
    ok = _mini_spec([{"type": "customers", "head": "H.", "caption": "C.",
                      "logo_set": "financial-services"}])
    assert compose.validate_spec(ok) == []
    bad = _mini_spec([{"type": "customers", "head": "H.", "caption": "C.",
                       "logo_set": "nope"}])
    assert any("unknown logo_set" in e for e in compose.validate_spec(bad))
    # a product-deck wall (industry cut) builds without error
    dest = compose.build(ok, tmp_path, dummy_assets)
    from pptx import Presentation
    assert len(Presentation(str(dest)).slides._sldIdLst) == 1


@pytest.mark.skipif(
    not (REPO / "workspace/clients/brisken/context/decks/reference-2026/"
         "OnePilot Solutions Overview 2026.pptx").exists(),
    reason="reference mirror not present (gitignored; CI)")
def test_badge_identity_real_reference():
    import hashlib
    ref = (REPO / "workspace/clients/brisken/context/decks/reference-2026/"
           "OnePilot Solutions Overview 2026.pptx")
    with zipfile.ZipFile(ref) as z:
        for src, (dst, pin) in assets_mod.BRAND_ASSETS.items():
            got = hashlib.md5(z.read(f"ppt/media/{src}")).hexdigest()
            assert got == pin, f"{src} -> {dst}: reference media drifted"


def test_accent_isolation(tmp_path, dummy_assets):
    """Each palette's build carries its own accent hex and no sibling
    product accent (the one-family / per-deck-accent invariant)."""
    body = [ALL_TYPES_SPEC[2], ALL_TYPES_SPEC[9], ALL_TYPES_SPEC[10]]
    hexes = {name: str(p.accent).upper() for name, p in T.PALETTES.items()}
    for name in T.PALETTES:
        out = tmp_path / name
        out.mkdir()
        dest = compose.build(_mini_spec(body, palette=name, output=f"iso-{name}"),
                             out, dummy_assets)
        with zipfile.ZipFile(dest) as z:
            xml = b"".join(z.read(n) for n in z.namelist()
                           if n.startswith("ppt/slides/") and n.endswith(".xml")).decode()
        assert hexes[name] in xml, f"{name}: own accent missing"
        for other, hx in hexes.items():
            if other != name:
                assert hx not in xml, f"{name}: contains {other}'s accent {hx}"
