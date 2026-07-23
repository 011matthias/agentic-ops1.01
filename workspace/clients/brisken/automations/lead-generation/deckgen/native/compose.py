# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "pillow", "pyyaml"]
# ///
"""Compose a native-family deck from its spec.

    uv run compose.py <deck>          # native/specs/<deck>.yaml -> dist pptx
    uv run compose.py <deck> --out D  # override output dir (tests)

Pipeline: load spec -> resolve palette -> ensure assets -> CLEAR the deck's
dist folder (v2 and native share dist/<deck>/; a stale v2 file must never
survive into an upload) -> build slides -> save -> gates (em-dash zero on
slide XML, banned-terms via tools/validate-demo-material.py).

Spec shape: {deck, output, palette, slides: [{type, ...content}]}. Slide
types map 1:1 onto grammar.Deck builders; compose owns pagination (cover /
divider / contact carry no page number) and fails loudly on overflow
(the use-case step-count collision class was fixed twice in the Overview
loop; here it is a build error, not a render surprise).
"""
from __future__ import annotations

import importlib.util
import subprocess
import sys
import zipfile
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Emu

_here = Path(__file__).parent


def _load(name, path=None):
    spec = importlib.util.spec_from_file_location(name, path or (_here / f"{name}.py"))
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


T = _load("tokens")
grammar = _load("grammar")
assets = _load("assets")
common = _load("common", _here.parent / "common.py")

UNNUMBERED = {"cover", "divider", "contact"}
DARK = {"divider", "contact"}

REQUIRED = {
    "cover": ["eyebrow", "title", "sub"],
    "divider": ["num", "title"],
    "short_version": ["head", "cards"],
    "about": ["kick", "head", "intro", "cards"],
    "hub": ["kick", "head", "intro", "cards"],
    "customers": ["head", "caption"],
    "governance": ["kick", "head", "intro", "cards", "footline"],
    "hierarchy": [],
    "functional": ["kick", "head", "product", "sources", "dests"],
    "app": ["kick", "head", "gloss", "problem", "steps", "freed", "connects"],
    "usecase": ["kick", "title", "gloss", "before", "steps", "checkpoint"],
    "problem": ["kick", "head", "bands"],
    "success": ["head", "intro", "stories"],
    "compare": ["kick", "head", "left_label", "left_items", "right_label", "right_items"],
    "contact": [],
}

# (list field, max entries) per type — fail loudly instead of rendering a collision.
LIMITS = {
    "usecase": [("steps", 6)],
    "app": [("steps", 5), ("problem", 4), ("freed", 4)],
    "problem": [("bands", 3)],
    "short_version": [("cards", 4)],
    "about": [("cards", 3)],
    "hub": [("cards", 3)],
    "governance": [("cards", 4)],
    "success": [("stories", 3)],
    "functional": [("sources", 5), ("dests", 5)],
    "compare": [("left_items", 4), ("right_items", 4)],
}


def _walk_strings(node):
    if isinstance(node, str):
        yield node
    elif isinstance(node, dict):
        for v in node.values():
            yield from _walk_strings(v)
    elif isinstance(node, list):
        for v in node:
            yield from _walk_strings(v)


def validate_spec(spec: dict) -> list[str]:
    errs = []
    for key in ("deck", "output", "palette", "slides"):
        if key not in spec:
            errs.append(f"spec missing top-level '{key}'")
    if errs:
        return errs
    if spec["palette"] not in T.PALETTES:
        errs.append(f"unknown palette '{spec['palette']}' (known: {sorted(T.PALETTES)})")
    for i, sl in enumerate(spec["slides"], 1):
        t = sl.get("type")
        if t not in REQUIRED:
            errs.append(f"slide {i}: unknown type '{t}' (known: {sorted(REQUIRED)})")
            continue
        for f in REQUIRED[t]:
            if f not in sl:
                errs.append(f"slide {i} ({t}): missing required field '{f}'")
        for field, cap in LIMITS.get(t, []):
            items = sl.get(field)
            if isinstance(items, list) and len(items) > cap:
                errs.append(f"slide {i} ({t}): {field} has {len(items)} entries, max {cap} "
                            f"(overflow collides with the layout — split the slide)")
        if t == "problem" and isinstance(sl.get("bands"), list) and len(sl["bands"]) != 3:
            errs.append(f"slide {i} (problem): needs exactly 3 bands")
    for s in _walk_strings(spec):
        if "—" in s:
            errs.append(f"em-dash in spec text: {s[:80]!r}")
    return errs


def build(spec: dict, out_dir: Path, assets_dir: Path) -> Path:
    pal = T.PALETTES[spec["palette"]]
    prs = Presentation()
    prs.slide_width = Emu(int(T.SW * T.EMU))
    prs.slide_height = Emu(int(T.SH * T.EMU))
    deck = grammar.Deck(prs, pal, assets_dir)
    page = 0
    for sl in spec["slides"]:
        t = sl["type"]
        kw = {k: v for k, v in sl.items() if k != "type"}
        if t not in UNNUMBERED:
            page += 1
            kw["page"] = page
        if t == "cover":
            s = deck.cover(**kw)
        elif t == "divider":
            s = deck.divider(**kw)
        elif t == "short_version":
            s = deck.short_version(**kw)
        elif t == "about":
            s = deck.about(**kw)
        elif t == "hub":
            s = deck.hub(**kw)
        elif t == "customers":
            s = deck.customers(**kw)
        elif t == "governance":
            s = deck.governance(**kw)
        elif t == "hierarchy":
            s = deck.hierarchy(**kw)
        elif t == "functional":
            s = deck.functional(**kw)
        elif t == "app":
            s = deck.app(**kw)
        elif t == "usecase":
            s = deck.usecase(**kw)
        elif t == "problem":
            s = deck.problem(**kw)
        elif t == "success":
            s = deck.success(**kw)
        elif t == "compare":
            s = deck.compare(**kw)
        elif t == "contact":
            s = deck.contact(**kw)
        if t not in DARK:
            deck.signature(s)
    dest = out_dir / f"{spec['output']}.pptx"
    prs.save(str(dest))
    return dest


def post_gates(pptx: Path) -> list[str]:
    errs = []
    with zipfile.ZipFile(pptx) as z:
        for n in z.namelist():
            if n.startswith("ppt/slides/") and n.endswith(".xml"):
                if "—".encode() in z.read(n) or b"\xe2\x80\x94" in z.read(n):
                    errs.append(f"em-dash in {n}")
    validator = common.main_clone_root() / "tools" / "validate-demo-material.py"
    res = subprocess.run(["uv", "run", str(validator), "--client", "brisken", str(pptx)],
                         capture_output=True, text=True)
    if res.returncode != 0:
        errs.append("banned-terms gate failed:\n" + res.stdout.strip() + res.stderr.strip())
    return errs


def main() -> int:
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if len(args) != 1:
        print(__doc__)
        return 2
    deck_name = args[0]
    spec_path = _here / "specs" / f"{deck_name}.yaml"
    if not spec_path.exists():
        print(f"no spec: {spec_path}")
        return 1
    spec = yaml.safe_load(spec_path.read_text(encoding="utf-8"))
    errs = validate_spec(spec)
    if errs:
        print("SPEC INVALID:")
        for e in errs:
            print(f"  - {e}")
        return 1

    if "--out" in sys.argv:
        out_dir = Path(sys.argv[sys.argv.index("--out") + 1])
        out_dir.mkdir(parents=True, exist_ok=True)
    else:
        out_dir = common.dist_dir() / deck_name
        if out_dir.exists():
            for f in out_dir.iterdir():
                f.unlink()
        out_dir.mkdir(parents=True, exist_ok=True)

    adir = assets.ensure()
    dest = build(spec, out_dir, adir)
    gate_errs = post_gates(dest)
    if gate_errs:
        print("GATES FAILED:")
        for e in gate_errs:
            print(f"  - {e}")
        return 1
    n = len(Presentation(str(dest)).slides._sldIdLst)
    print(f"composed: {dest}  ({n} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
