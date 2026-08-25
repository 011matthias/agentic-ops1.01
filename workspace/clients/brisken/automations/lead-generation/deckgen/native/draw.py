# /// script
# requires-python = ">=3.10"
# dependencies = ["python-pptx", "pillow"]
# ///
"""Drawing primitives for the native deck family. Pure geometry + text; no
deck content and no palette opinions live here (grammar.py owns those).
Ported verbatim from the approved Overview's build source so the shipped
deck regenerates identically."""
import importlib.util
import pathlib

from PIL import Image
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

_here = pathlib.Path(__file__).parent
_spec = importlib.util.spec_from_file_location("tokens", _here / "tokens.py")
T = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(T)


def IN(v):
    return Inches(v)


def _no_shadow(sh):
    sh.shadow.inherit = False


def _set_dash(line, val="dash"):
    ln = line._get_or_add_ln()
    for e in ln.findall(qn("a:prstDash")):
        ln.remove(e)
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": val}))


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False,
         radius=0.06, dash=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        IN(x), IN(y), IN(w), IN(h))
    _no_shadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
        if dash:
            _set_dash(shp.line, dash)
    return shp


def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {text, font, size, color, bold, align, space_before,
    space_after, line_spacing, spc(letterspacing pt)}."""
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if "space_before" in p:
            para.space_before = Pt(p["space_before"])
        if "space_after" in p:
            para.space_after = Pt(p["space_after"])
        if "line_spacing" in p:
            para.line_spacing = p["line_spacing"]
        run = para.add_run()
        run.text = p["text"]
        f = run.font
        f.name = p.get("font", T.BODY)
        f.size = Pt(p.get("size", 18))
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.color.rgb = p.get("color", T.INK)
        if "spc" in p:
            run.font._rPr.set("spc", str(int(p["spc"] * 100)))
    return tb


def place_image(slide, path, bx, by, bw, bh, halign="mid", valign="mid", max_h=None):
    """Aspect-fit an image within box (bx,by,bw,bh); center by default."""
    iw, ih = Image.open(path).size
    scale = min(bw / (iw / 96.0), bh / (ih / 96.0))
    w = (iw / 96.0) * scale
    h = (ih / 96.0) * scale
    if max_h and h > max_h:
        k = max_h / h
        w *= k
        h *= k
    x = bx + (bw - w) / 2 if halign == "mid" else (bx if halign == "left" else bx + bw - w)
    y = by + (bh - h) / 2 if valign == "mid" else (by if valign == "top" else by + bh - h)
    return slide.shapes.add_picture(str(path), IN(x), IN(y), IN(w), IN(h))


def blank(prs, bg=T.PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    if bg is not None:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
    return s


def arrow(slide, x, y, w, h, color):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, IN(x), IN(y), IN(w), IN(h))
    _no_shadow(a)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    try:
        a.adjustments[0] = 0.55
        a.adjustments[1] = 0.55
    except Exception:
        pass
    return a
