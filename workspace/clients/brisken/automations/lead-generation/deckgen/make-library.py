# /// script
# dependencies = ["python-pptx>=1.0.2"]
# ///
"""Derive the slide library from the pristine reference deck.

    uv run make-library.py

Deterministic + idempotent: always rebuilds library.pptx from the reference
copy in context (run fetch-reference.py first to refresh that). Steps:

1. Copy reference -> library.pptx (masters/layouts/theme/media/embedded
   fonts byte-identical).
2. Rename the patchable shapes of each donor slide to stable semantic role
   names (fails loudly if Dirk's restructuring moved a shape — that is the
   signal to update the mapping below against the new reference).
3. Author the three patterns the content architecture needs that the
   reference has no donor for — short-version, problem shell,
   platform-context hierarchy — using theme-native fonts/colors only.
4. Emit library-manifest.json (pattern -> donor slide -> roles).

Composition then COPIES library.pptx and duplicates/patches/deletes inside
the copy (see compose.py); this script never touches the reference.
"""

from __future__ import annotations

import importlib.util
import json
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

_here = Path(__file__).parent


def _load(name: str, fname: str):
    spec = importlib.util.spec_from_file_location(name, _here / fname)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


common = _load("common", "common.py")

# load the engine from the repo checkout that contains this file (worktree or
# main clone) so code+engine always version together
_engine_path = Path(__file__).resolve().parents[6] / "tools" / "pptx_slide_ops.py"
spec = importlib.util.spec_from_file_location("pptx_slide_ops", _engine_path)
ops = importlib.util.module_from_spec(spec)
sys.modules["pptx_slide_ops"] = ops
spec.loader.exec_module(ops)

SLATE = RGBColor(0x32, 0x3F, 0x4F)
CYAN = RGBColor(0x00, 0xB5, 0xD0)
LIGHT = RGBColor(0xD0, 0xD8, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# (slide_1based, current_name) -> semantic role name. First occurrence wins;
# every entry must resolve or the build fails (reference-drift tripwire).
RENAMES: list[tuple[int, str, str]] = [
    (1, "TextBox 19", "cover-title"),
    (1, "Picture 3", "cover-logo"),
    (3, "TextBox 33", "about-headline"),
    (3, "TextBox 34", "about-eyebrow"),
    # slide 6 has two shapes named 'Title 1'; first-occurrence renaming makes
    # the first entry claim the headline, the second the body
    (6, "Title 1", "hub-headline"),
    (6, "Title 1", "hub-body"),
    (7, "TextBox 121", "eco-headline"),
    (9, "Rounded Rectangle 1", "fdp-header"),
    (10, "TextBox 6", "transition-title"),
    (11, "TextBox 33", "vp-title"),
    (11, "TextBox 1", "vp-kicker"),
    (11, "Rectangle 17", "vp-card-main"),
    (11, "Rectangle 19", "vp-card-top"),
    (11, "Rectangle 18", "vp-card-bottom"),
    (11, "Picture Placeholder 37", "vp-hero"),
    (11, "Picture 13", "vp-logo"),
    (12, "Rectangle 1", "feat-body"),
    (12, "TextBox 3", "feat-header"),
    (12, "TextBox 6", "feat-logo-caption"),
    (13, "TextBox 18", "fd-header"),
    (13, "TextBox 1", "fd-gov-title"),
    (13, "TextBox 87", "fd-gov-list"),
    (13, "TextBox 8", "fd-src-1"),
    (13, "TextBox 35", "fd-src-2"),
    (13, "TextBox 72", "fd-src-3"),
    (13, "TextBox 13", "fd-src-4"),
    (13, "TextBox 57", "fd-src-5"),
    (13, "TextBox 19", "fd-dest-title"),
    (13, "TextBox 15", "fd-dest-list"),
    (13, "Title 7", "fd-title"),
    (15, "Rectangle 1", "feat-body"),
    (15, "TextBox 3", "feat-header"),
    (15, "TextBox 10", "feat-product-name"),
    (15, "TextBox 8", "feat-logo-caption"),
    (16, "TextBox 18", "fd-header"),
    (16, "TextBox 1", "fd-gov-title"),
    (16, "TextBox 87", "fd-gov-list"),
    (16, "TextBox 8", "fd-src-1"),
    (16, "TextBox 35", "fd-src-2"),
    (16, "TextBox 48", "fd-src-3"),
    (16, "TextBox 49", "fd-src-4"),
    (16, "TextBox 19", "fd-dest-title"),
    (16, "TextBox 15", "fd-dest-list"),
    (16, "Title 4", "fd-title"),
    (16, "TextBox 47", "fd-product-label"),
    (18, "Rectangle 1", "feat-body"),
    (18, "TextBox 3", "feat-header"),
    (18, "TextBox 5", "feat-product-name"),
    (18, "TextBox 6", "feat-logo-caption"),
    (19, "TextBox 18", "fd-header"),
    (19, "TextBox 1", "fd-gov-title"),
    (19, "TextBox 87", "fd-gov-list"),
    (19, "TextBox 8", "fd-src-1"),
    (19, "TextBox 35", "fd-src-2"),
    (19, "TextBox 72", "fd-src-3"),
    (19, "TextBox 13", "fd-src-4"),
    (19, "TextBox 57", "fd-src-5"),
    (19, "TextBox 19", "fd-dest-title"),
    (19, "TextBox 15", "fd-dest-list"),
    (19, "Rectangle 38", "fd-title"),
    (21, "Title 9", "uc-title"),
    (21, "Content Placeholder 12", "uc-col1-head"),
    (21, "Content Placeholder 11", "uc-col1-body"),
    (21, "Content Placeholder 14", "uc-col2-head"),
    (21, "Content Placeholder 13", "uc-col2-body"),
    (21, "TextBox 29", "uc-logo-name"),
    (21, "Picture 28", "uc-logo-image"),
    (28, "Title 17", "ss-title"),
    (28, "TextBox 19", "ss-sub"),
    (28, "TextBox 1", "ss-highlight"),
    (29, "Title 17", "ss-title"),
    (29, "TextBox 4", "ss-sub"),
    (29, "TextBox 5", "ss-highlight"),
    (30, "Title 17", "ss-title"),
    (30, "TextBox 3", "ss-sub"),
    (30, "TextBox 4", "ss-highlight"),
]

# pattern -> donor slide (1-based, in the LIBRARY = reference order + 3
# authored at the end) + the roles compose may address on it
MANIFEST: dict[str, dict] = {
    "cover":            {"donor": 1,  "roles": ["cover-title", "cover-logo"]},
    "about":            {"donor": 3,  "roles": ["about-headline", "about-eyebrow"]},
    "clients":          {"donor": 4,  "roles": []},
    "suite-hub":        {"donor": 6,  "roles": ["hub-headline", "hub-body"]},
    "ecosystem":        {"donor": 7,  "roles": ["eco-headline"]},
    "custom-dev":       {"donor": 8,  "roles": []},
    "func-diagram-platform": {"donor": 9, "roles": ["fdp-header"]},
    "transition":       {"donor": 10, "roles": ["transition-title"]},
    "value-3cards":     {"donor": 11, "roles": ["vp-title", "vp-kicker", "vp-card-main", "vp-card-top", "vp-card-bottom", "vp-hero", "vp-logo"]},
    "features":         {"donor": 12, "roles": ["feat-body", "feat-header", "feat-logo-caption"]},
    "func-diagram-mdh": {"donor": 13, "roles": ["fd-header", "fd-gov-title", "fd-gov-list", "fd-src-1", "fd-src-2", "fd-src-3", "fd-src-4", "fd-src-5", "fd-dest-title", "fd-dest-list", "fd-title"]},
    "features-st":      {"donor": 15, "roles": ["feat-body", "feat-header", "feat-product-name", "feat-logo-caption"]},
    "func-diagram-st":  {"donor": 16, "roles": ["fd-header", "fd-gov-title", "fd-gov-list", "fd-src-1", "fd-src-2", "fd-src-3", "fd-src-4", "fd-dest-title", "fd-dest-list", "fd-title", "fd-product-label"]},
    "features-dcw":     {"donor": 18, "roles": ["feat-body", "feat-header", "feat-product-name", "feat-logo-caption"]},
    "func-diagram-dcw": {"donor": 19, "roles": ["fd-header", "fd-gov-title", "fd-gov-list", "fd-src-1", "fd-src-2", "fd-src-3", "fd-src-4", "fd-src-5", "fd-dest-title", "fd-dest-list", "fd-title"]},
    "usecase":          {"donor": 21, "roles": ["uc-title", "uc-col1-head", "uc-col1-body", "uc-col2-head", "uc-col2-body", "uc-logo-name", "uc-logo-image"]},
    "usecase-cashflow": {"donor": 22, "roles": []},
    "usecase-bankfee":  {"donor": 23, "roles": []},
    "usecase-credit":   {"donor": 24, "roles": []},
    "usecase-esg":      {"donor": 25, "roles": []},
    "usecase-remit":    {"donor": 26, "roles": []},
    "success-fsi":      {"donor": 28, "roles": ["ss-title", "ss-sub", "ss-highlight"]},
    "success-agri":     {"donor": 29, "roles": ["ss-title", "ss-sub", "ss-highlight"]},
    "success-chem":     {"donor": 30, "roles": ["ss-title", "ss-sub", "ss-highlight"]},
    "contact":          {"donor": 31, "roles": []},
    "short-version":    {"donor": 32, "roles": ["sv-eyebrow", "sv-rows", "sv-handoff"], "authored": True},
    "problem":          {"donor": 33, "roles": ["pb-eyebrow", "pb-headline", "pb-body"], "authored": True},
    "platform-context": {"donor": 34, "roles": ["pc-title", "pc-onepilot", "pc-tc", "pc-app-mdh", "pc-app-st", "pc-app-dw", "pc-uc-strip", "pc-caption"], "authored": True,
                          "note": "compose highlights one pc-app-* box (cyan fill, white text)"},
}


def _txt(slide, name, text, x, y, w, h, size, bold=False, color=SLATE,
         font="Poppins", align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    r = tf.paragraphs[0].runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if align is not None:
        tf.paragraphs[0].alignment = align
    return tb


def _box(slide, name, text, x, y, w, h, fill, line, text_color, size, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.name = name
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.text = text
    for para in tf.paragraphs:
        for r in para.runs:
            r.font.name = "Poppins"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = text_color
    return shp


def author_short_version(prs, blank_layout):
    s = prs.slides.add_slide(blank_layout)
    _txt(s, "sv-eyebrow", "THE SHORT VERSION", 0.66, 0.70, 4.0, 0.35, 11)
    rows = _txt(s, "sv-rows", "What it is: ", 0.66, 1.60, 11.8, 4.2, 18)
    tf = rows.text_frame
    for extra in ("What it replaces: ", "What you get: ", "Who it is for: "):
        p = tf.add_paragraph()
        p.text = extra
        r = p.runs[0]
        r.font.name = "Poppins"
        r.font.size = Pt(18)
        r.font.color.rgb = SLATE
        p.space_after = Pt(14)
    tf.paragraphs[0].space_after = Pt(14)
    _txt(s, "sv-handoff", "The next page is where you are today.", 0.66, 6.30, 8.0, 0.45, 14, color=CYAN)
    return s


def author_problem(prs, blank_layout):
    s = prs.slides.add_slide(blank_layout)
    _txt(s, "pb-eyebrow", "THE PROBLEM", 0.66, 0.70, 4.0, 0.35, 11)
    _txt(s, "pb-headline", "Problem headline", 0.66, 1.30, 11.8, 1.3, 32, bold=True)
    body = _txt(s, "pb-body", "Problem body", 0.66, 2.90, 11.8, 3.9, 16)
    body.text_frame.paragraphs[0].space_after = Pt(12)
    return s


def author_platform_context(prs, blank_layout):
    s = prs.slides.add_slide(blank_layout)
    _txt(s, "pc-title", "WHERE IT SITS", 0.66, 0.70, 5.0, 0.35, 11)
    # platform band
    _box(s, "pc-onepilot",
         "OnePilot: the platform. Connected to everything, especially the SAP world. Customers build their own applications on it.",
         0.90, 5.05, 11.5, 1.05, SLATE, None, WHITE, 14)
    # TreasuryCentral band on the platform
    _box(s, "pc-tc", "TreasuryCentral: the screen your team works in",
         1.55, 3.95, 10.2, 0.85, WHITE, SLATE, SLATE, 13)
    # three applications inside TreasuryCentral
    _box(s, "pc-app-mdh", "Market Data Hub", 1.90, 2.30, 2.95, 1.30, LIGHT, None, SLATE, 14)
    _box(s, "pc-app-st", "Smart Trading", 5.15, 2.30, 2.95, 1.30, LIGHT, None, SLATE, 14)
    # brisken.com/onepilot names the AI application "AI Digital Workforce"
    _box(s, "pc-app-dw", "AI Digital Workforce", 8.40, 2.30, 2.95, 1.30, LIGHT, None, SLATE, 14)
    # use cases hang off the Digital Worker
    _box(s, "pc-uc-strip",
         "Use cases: funding requests, remittance advices, bank statements, cash positioning",
         6.55, 1.15, 6.0, 0.75, WHITE, CYAN, SLATE, 11, bold=False)
    _txt(s, "pc-caption", "Caption", 0.90, 6.45, 11.5, 0.6, 12)
    return s


def main() -> int:
    ref = common.reference_path()
    if not ref.exists():
        print(f"reference missing at {ref} — run fetch-reference.py first")
        return 1
    lib = common.library_path()
    shutil.copyfile(ref, lib)
    prs = Presentation(str(lib))
    n_ref = len(prs.slides)

    for slide_no, old, new in RENAMES:
        slide = prs.slides[slide_no - 1]
        shp = ops.find_shape(slide, old, required=False)
        if shp is None:
            # already renamed on a rebuild? semantic name present is fine
            if ops.find_shape(slide, new, required=False) is not None:
                continue
            raise SystemExit(
                f"RENAME FAIL slide {slide_no}: {old!r} not found — reference "
                f"structure drifted; update RENAMES against the new deck"
            )
        shp.name = new

    blank_layout = prs.slides[5].slide_layout  # the product sections' Blank
    author_short_version(prs, blank_layout)
    author_problem(prs, blank_layout)
    author_platform_context(prs, blank_layout)

    prs.save(str(lib))

    # verify: roles resolvable, fonts survived, rIds valid
    prs2 = Presentation(str(lib))
    assert len(prs2.slides) == n_ref + 3, "authored-slide count mismatch"
    for pat, info in MANIFEST.items():
        slide = prs2.slides[info["donor"] - 1]
        for role in info["roles"]:
            ops.find_shape(slide, role)  # raises if missing
    fonts = ops.check_embedded_fonts(lib)
    rids = ops.validate_rids(lib)
    if not fonts["ok"] or rids:
        print("fonts:", fonts, "rid problems:", rids)
        return 1

    manifest_path = _here / "library-manifest.json"
    manifest_path.write_text(
        json.dumps({"reference_slides": n_ref, "library_slides": n_ref + 3,
                    "patterns": MANIFEST}, indent=2) + "\n",
        encoding="utf-8",
    )
    print(f"library: {lib}  ({n_ref}+3 slides)")
    print(f"manifest: {manifest_path}")
    print(f"fonts: {fonts['font_parts']} parts OK; rIds OK; all roles resolve")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
